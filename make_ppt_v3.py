"""领导汇报PPT - 浅色对齐版"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def R(*c): return RGBColor(*c)

C_NAVY  = (26, 54, 93)
C_BLUE  = (0, 82, 170)
C_LIGHT = (235, 242, 255)
C_WHITE = (255, 255, 255)
C_GRAY1 = (248, 248, 248)
C_GRAY2 = (220, 220, 230)
C_TEXT  = (60, 60, 80)
C_SUB   = (120, 130, 160)
C_ACCENT= (230, 95, 0)

def rect(sl, l, t, w, h, fill, line=None, lw=0.75):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = R(*fill)
    if line:
        s.line.color.rgb = R(*line); s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def tb(sl, text, l, t, w, h, sz=9, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, wrap=True):
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold; run.font.color.rgb = R(*color)
    return box

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
sl = prs.slides.add_slide(prs.slide_layouts[6])

rect(sl, 0, 0, 13.333, 7.5, C_WHITE)

# 顶部横幅
rect(sl, 0, 0, 13.333, 0.85, C_NAVY)
rect(sl, 0, 0.8, 13.333, 0.05, C_ACCENT)
tb(sl, "ENTERPRISE ASSET & SECURITY MANAGEMENT", 0.5, 0.08, 10, 0.25, 8, False, (160,190,230)
tb(sl, "企业资产与信息安全综合管理平台", 0.5, 0.32, 10, 0.45, 20, True)
tb(sl, "B/S三层架构  |  Flask + MySQL + Ollama AI  |  v1.3", 0.5, 0.72, 12, 0.22, 9, False, (140,170,210)
tb(sl, "2026年5月8日", 11.2, 0.72, 2.0, 0.22, 9, False, (140,170,210, PP_ALIGN.RIGHT)

# ===== 第一行：三大成果 =====
y1 = 1.05
tb(sl, "核心成果", 0.5, y1, 2.0, 0.28, 9, True, C_NAVY)

cards = [
    ("数据集中化", "93,000+", "员工档案统一管理，消除信息孤岛"),
    ("权限精细化", "7级 RBAC", "最小粒度到人-系统-操作，部门数据隔离"),
    ("AI 智能化", "自然语言查数据", "非技术人员自主查询，秒级响应"),
]
cx = 0.5
for title, num, desc in cards:
    rect(sl, cx, y1 + 0.38, 4.1, 0.85, C_NAVY)
    tb(sl, title, cx + 0.2, y1 + 0.42, 1.5, 0.22, 9, True)
    tb(sl, num, cx + 0.2, y1 + 0.65, 3.7, 0.3, 16, True)
    tb(sl, desc, cx + 0.2, y1 + 0.96, 3.7, 0.22, 8, False, (180,210,240)
    cx += 4.25

# ===== 第二行：已完成 =====
y2 = y1 + 1.32
rect(sl, 0.5, y2, 12.3, 0.28, C_LIGHT, C_BLUE, 0.5)
tb(sl, "已上线运行  |  6大核心模块", 0.65, y2 + 0.03, 12, 0.22, 9, True, C_BLUE)

done = [
    ("人员管理系统", "9.3万员工档案 · 80+字段"),
    ("IT资产管理", "1000+办公电脑 · 500+工控机"),
    ("安全管控平台", "涉密人/物/区域/文件全管控"),
    ("AI智能问答", "意图识别 · SQL生成 · 自然语言"),
    ("权限安全加固", "CSRF · XSS · 密码策略"),
    ("界面统一优化", "小米设计风格一致体验"),
]
dx = 0.5
for name, stats in done:
    rect(sl, dx, y2 + 0.36, 2.0, 0.6, C_BLUE)
    tb(sl, name, dx + 0.15, y2 + 0.4, 1.7, 0.28, 9, True)
    tb(sl, stats, dx + 0.15, y2 + 0.68, 1.7, 0.22, 7.5, False, (200,220,255)
    dx += 2.07

# ===== 第三行：收尾中 =====
y3 = y2 + 1.08
rect(sl, 0.5, y3, 12.3, 0.28, C_GRAY1, C_GRAY2, 0.5)
tb(sl, "收尾中  |  安全测试 + 全面回归测试 + 性能优化", 0.65, y3 + 0.03, 12, 0.22, 9, True, C_TEXT)

progs = [
    ("日志审计收尾", "操作留痕，合规审计追溯", 90),
    ("安全测试收尾", "代码审计 · 渗透测试", 90),
    ("全面回归测试", "功能 · 性能 · 安全", 75),
    ("数据库优化", "索引 · 字段长度 · 导入效率", 60),
    ("文档完善", "用户手册 · 部署文档", 50),
]
px = 0.5
for name, desc, pct in progs:
    pw = 2.35 * pct / 100
    rect(sl, px, y3 + 0.36, 2.35, 0.5, C_GRAY1, C_GRAY2, 0.5)
    rect(sl, px, y3 + 0.36, pw, 0.5, C_NAVY)
    tb(sl, f"{pct}%", px + 1.7, y3 + 0.4, 0.6, 0.28, 9, True, C_NAVY, PP_ALIGN.RIGHT)
    tb(sl, name, px + 0.15, y3 + 0.4, 1.5, 0.22, 9, True, C_TEXT)
    tb(sl, desc, px + 0.15, y3 + 0.62, 2.15, 0.2, 7.5, False, C_SUB)
    px += 2.45

# ===== 第四行：里程碑 =====
y4 = y3 + 1.0
rect(sl, 0.5, y4, 12.3, 0.28, C_LIGHT, C_BLUE, 0.5)
tb(sl, "里程碑  |  上线倒计时", 0.65, y4 + 0.03, 5, 0.22, 9, True, C_BLUE)

ms = [
    ("5/14", "安全收尾完成"),
    ("5/21", "全面测试完成"),
    ("5/28", "优化文档完成"),
    ("6/7", "正式上线运行"),
]
mx = 0.7
for date, title in ms:
    bg_c = C_ACCENT if "6/7" in date else C_NAVY
    rect(sl, mx, y4 + 0.36, 2.85, 0.5, bg_c)
    tb(sl, date, mx + 0.15, y4 + 0.38, 1.0, 0.28, 12, True)
    tb(sl, title, mx + 0.15, y4 + 0.66, 2.55, 0.2, 9, True, (220,230,255)
    mx += 3.05

# 里程碑箭头
for ax in [3.65, 6.7, 9.75]:
    rect(sl, ax, y4 + 0.58, 0.2, 0.04, C_GRAY2)
    tb(sl, "▶", ax - 0.05, y4 + 0.5, 0.3, 0.2, 10, False, C_GRAY2)

# ===== 底部统计栏 =====
rect(sl, 0, 5.62, 13.333, 1.35, C_NAVY)
tb(sl, "系统规模", 0.5, 5.68, 1.5, 0.25, 9, True, (180,200,230)

stats = [
    ("93,000+", "员工档案"),
    ("1,000+", "办公电脑"),
    ("500+", "工控机"),
    ("200+", "部门"),
    ("7级", "角色权限"),
    ("CSRF+XSS+SQL", "安全防护"),
    ("AI", "智能问答"),
    ("50+", "API接口"),
    ("秒级", "查询响应"),
]
sx = 0.5
for num, unit in stats:
    rect(sl, sx, 5.98, 1.2, 0.75, (40,70,110)
    tb(sl, num, sx + 0.08, 6.02, 1.04, 11, True)
    tb(sl, unit, sx + 0.08, 6.34, 1.04, 8, False, (150,180,220), PP_ALIGN.CENTER)
    sx += 1.26

# 预计上线
rect(sl, 12.0, 5.62, 1.2, 1.35, C_ACCENT)
tb(sl, "6/7", 12.05, 5.75, 1.0, 0.4, 18, True)
tb(sl, "预计上线", 12.05, 6.1, 1.0, 0.25, 9, True, (255,220,180))
tb(sl, "Enterprise Asset & Security Platform v1.3  |  2026-05-08", 0.5, 6.88, 12, 0.2, 8, False, (130,150,180), PP_ALIGN.CENTER)

prs.save("d:\\资产管理\\领导汇报_浅色对齐版.pptx")
print("OK")
