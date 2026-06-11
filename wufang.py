from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY=(26,54,93);ACCENT=(230,95,0);WHITE=(255,255,255)
GRAY=(248,248,248);TEXT=(60,60,80)
GREEN=(52,211,153);BLUE=(0,82,170);RED=(200,60,60)
PURPLE=(112,48,160);GRAY2=(220,220,235)

def R(r,g,b): return RGBColor(r,g,b)

def rect(sl,l,t,w,h,fill,lc=None,lw=0.75):
    s=sl.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid();s.fill.fore_color.rgb=R(*fill)
    if lc: s.line.color.rgb=R(*lc);s.line.width=Pt(lw)
    else: s.line.fill.background()
    return s

def tb(sl,txt,l,t,w,h,sz=9,bold=False,color=TEXT,align=PP_ALIGN.LEFT):
    box=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=box.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.alignment=align
    run=p.add_run();run.text=txt;run.font.size=Pt(sz);run.font.bold=bold
    run.font.color.rgb=R(*color)
    return box

prs=Presentation()
prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5)
sl=prs.slides.add_slide(prs.slide_layouts[6])

rect(sl,0,0,13.333,7.5,WHITE)
rect(sl,0,0,13.333,0.85,NAVY)
rect(sl,0,0.8,13.333,0.05,ACCENT)
tb(sl,"ASSET & SECURITY MANAGEMENT",0.5,0.08,10,0.25,8,False,(160,190,230))
tb(sl,"资产管理系统如何支撑企业\u201c五防\u201d管控",0.5,0.32,12,0.45,20,True)
tb(sl,"2026年5月8日",11.2,0.72,2.0,0.22,9,False,(140,170,210),PP_ALIGN.RIGHT)

items=[
    ("防越权","部门数据隔离\n7级RBAC权限\nAPI Key+CSRF认证",ACCENT),
    ("防泄露","涉密人/物/区域管控\n导出留痕+字段保护\n涉密文件全生命周期",BLUE),
    ("防丢失","人机绑定可追溯\n资产全生命周期管理\n介质责任到人",GREEN),
    ("防违规","操作全程留痕\n邮件服务器配置\n密码策略统一",RED),
    ("防篡改","权限最小化\n修改有记录\nCSRF+XSS+SQL防护",PURPLE),
]
ix=0.5
for name,desc,c in items:
    rect(sl,ix,1.1,2.4,2.9,WHITE,c,2)
    rect(sl,ix,1.1,2.4,0.55,c)
    tb(sl,name,ix+0.1,1.15,2.2,0.45,14,True)
    rect(sl,ix+0.1,1.75,2.2,0.01,GRAY)
    tb(sl,desc,ix+0.15,1.9,2.1,2.0,8,False,TEXT)
    ix+=2.55

tb(sl,"五防管控覆盖模块",0.5,4.15,3,0.25,9,True,NAVY)
headers=["维度","人员管理","资产管理","安全管理","AI服务","系统管理"]
rows=[
    ["防越权","部门隔离","权限过滤","角色授权","API认证","7级RBAC"],
    ["防泄露","字段保护","导出留痕","涉密管控","数据脱敏","日志审计"],
    ["防丢失","档案追踪","人机绑定","责任到人","操作日志","资产盘点"],
    ["防违规","密码策略","流程审批","涉密规范","请求审计","操作留痕"],
    ["防篡改","修改记录","删除审核","CSRF防护","输入校验","最小权限"],
]
for i,h in enumerate(headers):
    w=2.05 if i==0 else 2.0
    x=0.5+sum([2.05 if j==0 else 2.0 for j in range(i)])
    rect(sl,x,4.5,w,0.35,NAVY)
    tb(sl,h,x+0.08,4.52,w-0.16,0.3,8,True)
for ri,row in enumerate(rows):
    y=4.9+ri*0.32;bg_c=GRAY if ri%2==0 else WHITE
    for ci,cell in enumerate(row):
        w=2.05 if ci==0 else 2.0
        x=0.5+sum([2.05 if j==0 else 2.0 for j in range(ci)])
        rect(sl,x,y,w,0.32,bg_c,GRAY2,0.5)
        tb(sl,cell,x+0.08,y+0.04,w-0.16,0.28,7,ci==0,NAVY if ci==0 else TEXT)

tb(sl,"Enterprise Asset & Security Platform v1.3 | 2026-05-08",0.5,6.72,12,0.2,8,False,(130,150,180))
prs.save(r"d:\资产管理\五防管控.pptx")
print("done")
