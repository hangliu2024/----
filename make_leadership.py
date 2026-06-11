"""领导汇报PPT - 单页版"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def rgb(*c): return RGBColor(*c)

def bg(slide, l, t, w, h, c, lc=None, lw=0.5):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(*c)
    if lc: s.line.color.rgb = rgb(*lc); s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, sz, bold=False, color=(255,255,255), align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.color.rgb = rgb(*color)
    return tb

def mtxt(slide, lines, l, t, w, h, sizes, colors, bolds, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i,(line,sz,col,bold) in enumerate(zip(lines,sizes,colors,bolds)):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.size = Pt(sz); run.font.bold = bold
        run.font.color.rgb = rgb(*col)
    return tb

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])

# 白色背景
bg(s, 0, 0, 13.333, 7.5, (255,255,255))

# 顶部深蓝标题栏
bg(s, 0, 0, 13.333, 0.9, (26,54,93))
bg(s, 0, 0.83, 13.333, 0.07, (255,103,0))
txt(s, "ENTERPRISE ASSET & SECURITY MANAGEMENT", 0.4, 0.06, 8, 0.28, 8, False, (160,190,230))
txt(s, "企业资产与信息安全综合管理平台", 0.4, 0.32, 10, 0.5, 20, True)
txt(s, "B/S三层架构  |  Flask + MySQL + Ollama AI  |  v1.3", 0.4, 0.72, 12, 0.22, 8, False, (140,170,210))
txt(s, "汇报日期：2026年5月8日", 10.8, 0.35, 2.3, 0.3, 9, False, (140,170,210), PP_ALIGN.RIGHT)

# ===== 左列：核心业务价值 =====
bg(s, 0.4, 1.05, 5.8, 0.3, (26,54,93))
txt(s, "核心业务价值", 0.5, 1.07, 5.5, 0.26, 10, True)

values = [
    ("数据集中化", "9.3万员工档案统一管理，消除信息孤岛", (255,103,0)),
    ("查询自助化", "自然语言查数据，响应周期从天缩短到秒", (0,112,210)),
    ("权限精细化", "最小粒度到人-系统-操作，部门数据隔离", (200,60,60)),
    ("资产可视化", "人机关联，归属清晰，丢失可快速定位", (15,140,80)),
    ("安全合规化", "操作留痕，权限到人，满足保密资质要求", (112,48,160)),
    ("AI智能化", "非技术人员自主查询，解放IT生产力", (0,128,128)),
]
y = 1.42
for title, detail, c in values:
    bg(s, 0.4, y, 0.1, 0.44, c)
    bg(s, 0.55, y, 5.65, 0.44, (247,250,252), lw=1)
    txt(s, title, 0.65, y+0.04, 2.0, 0.36, 10, True, (26,54,93))
    txt(s, detail, 0.65, y+0.22, 5.5, 0.2, 8, False, (80,80,80))
    y += 0.5

# ===== 中列：项目进展 =====
bg(s, 6.65, 1.05, 3.3, 0.3, (26,54,93))
txt(s, "项目进展", 6.75, 1.07, 3.0, 0.26, 10, True)

done = [
    ("人员管理系统", "93,000+员工档案", (52,211,153)),
    ("IT资产管理系统", "1,000+办公电脑+500+工控机", (0,112,210)),
    ("安全管控平台", "涉密人/物/区域/文件", (200,60,60)),
    ("AI智能问答", "自然语言秒级查数据", (112,48,160)),
    ("权限安全加固", "CSRF+开放重定向+密码策略", (255,103,0)),
    ("界面统一优化", "小米设计风格一致体验", (0,128,128)),
]
yd = 1.42
for name, desc, c in done:
    bg(s, 6.65, yd, 0.1, 0.44, c)
    bg(s, 6.8, yd, 3.15, 0.44, (247,250,252), lw=1)
    txt(s, name, 6.9, yd+0.04, 2.0, 0.2, 10, True, c)
    txt(s, desc, 6.9, yd+0.22, 3.0, 0.2, 7.5, False, (80,80,80))
    yd += 0.5

# ===== 右列：里程碑 =====
bg(s, 10.2, 1.05, 2.9, 0.3, (26,54,93))
txt(s, "里程碑节点", 10.3, 1.07, 2.7, 0.26, 10, True)

ms = [
    ("5/14", "安全收尾完成", (52,211,153)),
    ("5/21", "全面回归测试", (0,112,210)),
    ("5/28", "性能优化文档", (0,112,210)),
    ("6/7", "正式上线运行", (255,103,0)),
]
ym = 1.42
for date, title, c in ms:
    bg(s, 10.2, ym, 2.9, 0.44, (247,250,252), lw=1)
    bg(s, 10.2, ym, 0.1, 0.44, c)
    txt(s, date, 10.35, ym+0.04, 1.0, 0.2, 12, True, c)
    txt(s, title, 10.35, ym+0.24, 2.7, 0.18, 8, False, (60,60,80))
    ym += 0.5

# ===== 底部总结栏 =====
bg(s, 0.4, 4.62, 12.5, 0.02, (220,220,230))
bg(s, 0.4, 4.72, 12.5, 1.65, (26,54,93))

# 左侧统计数据
stats = [
    ("93,000+", "员工档案", (255,103,0)),
    ("1,000+", "办公电脑", (0,112,210)),
    ("500+", "工控机", (15,140,80)),
    ("7级", "角色权限", (200,60,60)),
    ("CSRF+XSS+SQL", "安全防护", (112,48,160)),
    ("50+", "API接口", (0,128,128)),
    ("秒级", "查询响应", (80,80,80)),
]
xs = 0.5
for num, unit, c in stats:
    bg(s, xs, 4.8, 1.65, 0.58, (40,70,110))
    txt(s, num, xs+0.05, 4.82, 1.55, 0.28, 11, True)
    txt(s, unit, xs+0.05, 5.1, 1.55, 0.24, 8, False, (150,180,220))
    xs += 1.73

# 右侧结论
bg(s, 12.5, 4.72, 0.35, 1.65, (255,103,0))
txt(s, "6/7", 12.55, 4.82, 0.35, 0.28, 14, True, (255,255,255))
txt(s, "预计上线", 12.55, 5.1, 0.35, 0.22, 8, False, (255,200,150), PP_ALIGN.CENTER)

# 版本信息
txt(s, "Enterprise Asset & Security Management Platform v1.3  |  数据截至：2026-05-08", 0.5, 6.45, 12.5, 0.22, 8, False, (130,150,180), PP_ALIGN.CENTER)

prs.save("d:\\资产管理\\领导汇报_2026年5月.pptx")
print("OK")
