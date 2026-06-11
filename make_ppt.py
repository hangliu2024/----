"""生成项目进展PPT - 领导汇报版"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

def add_text_box(slide, text, left, top, width, height, font_size, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    shape.line.fill.background()
    return shape

bg = add_rect(slide, 0, 0, 13.333, 7.5, (26, 54, 93))
header_bg = add_rect(slide, 0, 0, 13.333, 1.1, (255, 103, 0))

add_text_box(slide, "📊 项目进展汇报", 0.4, 0.25, 12, 0.6, 28, True, (255, 255, 255), PP_ALIGN.LEFT)
add_text_box(slide, f"汇报日期：{datetime.date.today().strftime('%Y年%m月%d日')}", 9.5, 0.25, 3.5, 0.6, 13, False, (255, 255, 255), PP_ALIGN.RIGHT)

sections = [
    ("✅ 人员管理系统", "93,000+ 员工档案，80+ 字段覆盖全生命周期", 100, (52, 211, 153)),
    ("✅ IT资产管理系统", "1,000+ 办公电脑 · 500+ 工控机，全生命周期管理", 100, (52, 211, 153)),
    ("✅ 安全管控平台", "涉密人/物/区域/文件四位一体，RBAC权限体系", 100, (52, 211, 153)),
    ("✅ AI智能问答", "自然语言查数据，秒级响应，无需SQL技能", 100, (52, 211, 153)),
    ("✅ 权限安全加固", "CSRF防护、部门数据隔离、密码策略统一", 100, (52, 211, 153)),
    ("✅ 界面统一优化", "小米设计风格，体验一致，交互友好", 100, (52, 211, 153)),
    ("🟡 日志审计收尾", "操作留痕，支持合规审计追溯", 90, (255, 193, 7)),
    ("🟡 安全测试收尾", "代码审计+渗透测试，6月前完成", 90, (255, 193, 7)),
]

y = 1.35
for name, desc, pct, color in sections:
    bar_bg = add_rect(slide, 0.4, y, 7.5, 0.45, (40, 40, 60))
    bar_fill = add_rect(slide, 0.4, y, 7.5 * pct / 100, 0.45, color)
    add_text_box(slide, name, 0.5, y + 0.03, 3.5, 0.4, 13, True, (255, 255, 255), PP_ALIGN.LEFT)
    add_text_box(slide, f"{pct}%", 7.0, y + 0.03, 0.8, 0.4, 12, True, (255, 255, 255), PP_ALIGN.RIGHT)
    add_text_box(slide, desc, 8.2, y + 0.05, 4.8, 0.35, 11, False, (200, 200, 220), PP_ALIGN.LEFT)
    y += 0.62

y2 = 1.35
milestones = [
    ("5/8", "安全加固 + 权限体系完成", "🔴 P0全部修复"),
    ("5/14", "安全收尾 + 第一轮测试", "🔴 P0收尾"),
    ("5/21", "全面回归测试", "🟠 P1功能验证"),
    ("5/28", "性能优化 + 文档完善", "🟡 P2优化"),
    ("6/7", "✅ 正式上线运行", "🎯 里程碑"),
]
mx = 9.4
for date, desc, tag in milestones:
    add_rect(slide, mx, y2, 1.35, 0.58, (60, 60, 90))
    add_text_box(slide, date, mx + 0.05, y2 + 0.02, 0.55, 0.28, 12, True, (255, 103, 0), PP_ALIGN.LEFT)
    add_text_box(slide, desc, mx + 0.05, y2 + 0.28, 1.25, 0.28, 8, False, (180, 180, 200), PP_ALIGN.LEFT)
    y2 += 0.65

add_rect(slide, 0.4, 6.4, 12.5, 0.02, (255, 255, 255))

overdue = add_rect(slide, 0.4, 6.55, 3.5, 0.6, (52, 211, 153))
add_text_box(slide, "✅ 功能开发完成度：100%", 0.5, 6.6, 3.3, 0.5, 13, True, (26, 54, 93), PP_ALIGN.LEFT)

overdue2 = add_rect(slide, 4.2, 6.55, 3.5, 0.6, (255, 193, 7))
add_text_box(slide, "🟡 安全测试收尾：90%", 4.3, 6.6, 3.3, 0.5, 13, True, (26, 54, 93), PP_ALIGN.LEFT)

overdue3 = add_rect(slide, 8.0, 6.55, 4.5, 0.6, (255, 103, 0))
add_text_box(slide, "🎯 预计上线：2026年6月7日", 8.1, 6.6, 4.3, 0.5, 14, True, (255, 255, 255), PP_ALIGN.LEFT)

output_path = "d:\\资产管理\\项目进展汇报.pptx"
prs.save(output_path)
print(f"PPT已保存: {output_path}")