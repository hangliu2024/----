"""生成专业企业架构图PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import math

# 品牌色
C_NAVY    = (26, 54, 93)     # 深蓝主色
C_ORANGE  = (255, 103, 0)    # 橙色强调
C_BLUE    = (0, 112, 210)    # 微软蓝
C_GREEN   = (0, 176, 80)     # 绿色
C_RED     = (200, 50, 50)    # 红色
C_PURPLE  = (112, 48, 160)   # 紫色
C_TEAL    = (0, 128, 128)   # 青色
C_GRAY    = (89, 89, 89)     # 灰色
C_LGRAY   = (230, 230, 230)  # 浅灰
C_WHITE   = (255, 255, 255)
C_LBLUE   = (217, 226, 243)  # 浅蓝背景

def rgb(*args):
    return RGBColor(*args)

def add_rect(slide, l, t, w, h, fill, line=None, lw=0.75):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(*fill)
    if line:
        s.line.color.rgb = rgb(*line); s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def add_text(slide, text, l, t, w, h, size, bold=False, color=C_WHITE,
             align=PP_ALIGN.CENTER, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(*color)
    return tb

def add_icon_box(slide, icon, label, l, t, w, h, bg, border):
    """带图标的模块卡片"""
    # 背景
    r = add_rect(slide, l, t, w, h, bg, border, 1.5)
    # 顶部色条
    bar = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(*border)
    bar.line.fill.background()
    # 图标
    add_text(slide, icon, l + 0.1, t + 0.08, w - 0.2, 0.28, 14, False, border, PP_ALIGN.CENTER, False)
    # 标签
    add_text(slide, label, l + 0.05, t + 0.36, w - 0.1, h - 0.42, 8.5, False, C_NAVY, PP_ALIGN.CENTER, True)
    return r

def add_connector_v(slide, l, t, h, color=(150, 150, 170)):
    """垂直连接线"""
    r = add_rect(slide, l - 0.025, t, 0.05, h, color, None)

def add_connector_h(slide, l, t, w, color=(150, 150, 170)):
    """水平连接线"""
    r = add_rect(slide, l, t - 0.025, w, 0.05, color, None)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ===== 白色背景 =====
bg = add_rect(slide, 0, 0, 13.333, 7.5, C_WHITE, None)

# ===== 顶部标题栏 =====
title_bg = add_rect(slide, 0, 0, 13.333, 0.9, C_NAVY, None)
# 左侧标签
add_rect(slide, 0, 0.25, 0.12, 0.4, C_ORANGE, None)
add_text(slide, "ENTERPRISE ASSET & SECURITY MANAGEMENT", 0.2, 0.1, 5, 0.35, 8, False, (180, 200, 230), PP_ALIGN.LEFT, False)
add_text(slide, "企业资产与信息安全综合管理平台", 0.2, 0.42, 10, 0.4, 17, True, C_WHITE, PP_ALIGN.LEFT, False)
add_text(slide, "系统架构设计图  |  B/S 三层架构  |  技术栈: Flask + MySQL + Ollama", 0.2, 0.72, 12, 0.22, 8, False, (140, 170, 210), PP_ALIGN.LEFT, False)
add_text(slide, "v1.3  ·  2026-05-08", 11.0, 0.62, 2.0, 0.22, 8, False, (140, 170, 210), PP_ALIGN.RIGHT, False)

# ===== 分层背景区域 =====

# 用户层
user_bg = add_rect(slide, 0.4, 1.08, 12.5, 0.72, C_LBLUE, C_BLUE, 0.5)
add_text(slide, "PRESENTATION LAYER  ·  表现层", 0.5, 1.1, 2.8, 0.22, 8, True, C_BLUE, PP_ALIGN.LEFT, False)
add_text(slide, "浏览器 / 移动端 / API 调用", 3.4, 1.12, 3, 0.2, 7.5, False, (100, 130, 180), PP_ALIGN.LEFT, False)

# 网关层
gw_bg = add_rect(slide, 0.4, 1.88, 12.5, 0.52, (255, 248, 235), C_ORANGE, 0.5)
add_text(slide, "SECURITY GATEWAY  ·  安全网关层", 0.5, 1.9, 2.8, 0.22, 8, True, C_ORANGE, PP_ALIGN.LEFT, False)

# 服务层
svc_bg = add_rect(slide, 0.4, 2.48, 12.5, 0.85, (242, 247, 252), (0, 80, 160), 0.5)
add_text(slide, "APPLICATION LAYER  ·  应用服务层", 0.5, 2.5, 3.0, 0.22, 8, True, (0, 80, 160), PP_ALIGN.LEFT, False)

# 数据层
data_bg = add_rect(slide, 0.4, 3.42, 12.5, 0.85, (240, 250, 240), (0, 140, 80), 0.5)
add_text(slide, "DATA LAYER  ·  数据服务层", 0.5, 3.44, 3.0, 0.22, 8, True, (0, 140, 80), PP_ALIGN.LEFT, False)

# ===== 用户层模块 =====
users = [
    ("👑", "系统管理员", C_ORANGE),
    ("🏢", "部门主管", C_BLUE),
    ("👤", "普通员工", (80, 80, 100)),
    ("🔌", "API调用方", C_PURPLE),
]
ux = 0.7
for icon, name, c in users:
    card = slide.shapes.add_shape(1, Inches(ux), Inches(1.2), Inches(1.4), Inches(0.52))
    card.fill.solid(); card.fill.fore_color.rgb = rgb(*c)
    card.line.fill.background()
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{icon}\n{name}"
    run.font.size = Pt(8.5)
    run.font.bold = True
    run.font.color.rgb = rgb(*C_WHITE)
    ux += 1.5

# 用户层连接线
add_connector_v(slide, 6.2, 1.78, 0.12)

# ===== 网关层模块 =====
gw_items = [
    ("🛡️", "CSRF Token", "全局CSRF防护"),
    ("🔑", "API Key", "外部API认证"),
    ("👤", "JWT Session", "会话管理"),
    ("⚡", "Rate Limit", "限流熔断"),
]
gx = 0.7
for icon, name, desc in gw_items:
    box = slide.shapes.add_shape(1, Inches(gx), Inches(2.0), Inches(1.4), Inches(0.32))
    box.fill.solid(); box.fill.fore_color.rgb = rgb(255, 248, 235)
    box.line.color.rgb = rgb(255, 103, 0); box.line.width = Pt(0.75)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{icon} {name}"
    run.font.size = Pt(7.5)
    run.font.bold = True
    run.font.color.rgb = rgb(180, 80, 0)
    gx += 1.5

add_text(slide, "+ 更多安全策略", 0.7, 2.34, 1.4, 0.2, 7, False, (180, 100, 0), PP_ALIGN.CENTER, False)
add_connector_v(slide, 6.2, 2.42, 0.1)

# ===== 服务层模块 =====
services = [
    ("🔐", "认证服务", "/auth", C_NAVY),
    ("👥", "人员管理", "/personnel", C_BLUE),
    ("💻", "IT资产", "/assets", (15, 140, 80)),
    ("🛡️", "安全管控", "/security", C_RED),
    ("🤖", "AI助手", "/ai_assistant", C_PURPLE),
    ("⚙️", "系统管理", "/admin", (80, 80, 80)),
    ("📋", "日志审计", "/audit", (0, 128, 128)),
]
sx = 0.7
for icon, name, route, c in services:
    card = slide.shapes.add_shape(1, Inches(sx), Inches(2.62), Inches(1.55), Inches(0.62))
    card.fill.solid(); card.fill.fore_color.rgb = rgb(*c)
    card.line.fill.background()
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{icon}\n{name}"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = rgb(*C_WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = route
    run2.font.size = Pt(7)
    run2.font.color.rgb = rgb(200, 220, 255)
    sx += 1.7

add_connector_v(slide, 6.2, 3.3, 0.12)

# ===== 数据层模块 =====
# MySQL
db_card = slide.shapes.add_shape(1, Inches(0.7), Inches(3.58), Inches(3.6), Inches(0.62))
db_card.fill.solid(); db_card.fill.fore_color.rgb = rgb(230, 242, 255)
db_card.line.color.rgb = rgb(0, 80, 160); db_card.line.width = Pt(1.5)
tf = db_card.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "🗄  MySQL 8.0  数据库引擎"
run.font.size = Pt(11)
run.font.bold = True
run.font.color.rgb = rgb(0, 60, 130)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "11张核心业务表 · 主从复制 · 读写分离"
run2.font.size = Pt(7.5)
run2.font.color.rgb = rgb(80, 120, 180)
p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = "员工档案 · 资产台账 · 权限矩阵 · 操作日志"
run3.font.size = Pt(7)
run3.font.color.rgb = rgb(120, 150, 200)

# Redis
redis_card = slide.shapes.add_shape(1, Inches(4.6), Inches(3.58), Inches(3.0), Inches(0.62))
redis_card.fill.solid(); redis_card.fill.fore_color.rgb = rgb(255, 240, 248)
redis_card.line.color.rgb = rgb(180, 40, 100); redis_card.line.width = Pt(1.5)
tf = redis_card.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "⚡  Redis 分布式缓存"
run.font.size = Pt(11)
run.font.bold = True
run.font.color.rgb = rgb(160, 30, 80)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "会话存储 · 权限缓存 · 请求限流"
run2.font.size = Pt(7.5)
run2.font.color.rgb = rgb(180, 80, 130)

# Ollama
ai_card = slide.shapes.add_shape(1, Inches(7.9), Inches(3.58), Inches(3.8), Inches(0.62))
ai_card.fill.solid(); ai_card.fill.fore_color.rgb = rgb(245, 238, 255)
ai_card.line.color.rgb = rgb(130, 50, 180); ai_card.line.width = Pt(1.5)
tf = ai_card.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "🤖  Ollama AI 引擎 (Qwen3.5:9B)"
run.font.size = Pt(11)
run.font.bold = True
run.font.color.rgb = rgb(100, 30, 160)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "意图识别 · SQL生成 · 自然语言回答"
run2.font.size = Pt(7.5)
run2.font.color.rgb = rgb(140, 80, 200)

# ===== 右侧说明栏 =====
side_bg = add_rect(slide, 11.9, 1.08, 1.4, 3.2, (248, 248, 250), C_LGRAY, 0.5)
add_text(slide, "技术栈", 11.95, 1.12, 1.3, 0.22, 8, True, C_NAVY, PP_ALIGN.CENTER, False)

tech_stack = [
    ("Web", "Flask 2.x", C_NAVY),
    ("ORM", "SQLAlchemy", C_BLUE),
    ("DB", "MySQL 8.0", (0, 120, 80)),
    ("Cache", "Redis", C_RED),
    ("AI", "Ollama", C_PURPLE),
    ("Auth", "Flask-Login", C_TEAL),
    ("API", "RESTful", (80, 80, 80)),
    ("Sec", "CSRF+JWT", C_ORANGE),
]
ty = 1.4
for name, val, c in tech_stack:
    tag = slide.shapes.add_shape(1, Inches(11.95), Inches(ty), Inches(0.55), Inches(0.26))
    tag.fill.solid(); tag.fill.fore_color.rgb = rgb(*c)
    tag.line.fill.background()
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = name
    run.font.size = Pt(6.5)
    run.font.bold = True
    run.font.color.rgb = rgb(*C_WHITE)

    val_tb = slide.shapes.add_textbox(Inches(12.53), Inches(ty + 0.03), Inches(0.75), Inches(0.22))
    tf = val_tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = val
    run.font.size = Pt(7)
    run.font.color.rgb = rgb(*C_GRAY)
    ty += 0.37

# ===== 底部统计栏 =====
footer = add_rect(slide, 0, 4.45, 13.333, 1.2, (30, 60, 100), None)
add_text(slide, "系统规模与能力", 0.5, 4.5, 2.2, 0.28, 9, True, (180, 200, 230), PP_ALIGN.LEFT, False)

stats = [
    ("👥", "93,000+", "员工档案"),
    ("💻", "1,000+", "办公电脑"),
    ("🔧", "500+", "工控机"),
    ("🏢", "200+", "部门"),
    ("🔐", "7级", "角色权限"),
    ("🛡️", "CSRF+XSS+SQL", "全防护"),
    ("🤖", "AI", "智能问答"),
    ("📊", "50+", "API接口"),
    ("📱", "B/S+C/S", "混合架构"),
    ("⚡", "秒级", "查询响应"),
]
sx2 = 0.5; sy = 4.85
for icon, num, unit in stats:
    box = slide.shapes.add_shape(1, Inches(sx2), Inches(sy), Inches(1.15), Inches(0.6))
    box.fill.solid(); box.fill.fore_color.rgb = rgb(40, 70, 110)
    box.line.fill.background()
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{icon} {num}"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = rgb(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = unit
    run2.font.size = Pt(7)
    run2.font.color.rgb = rgb(150, 180, 220)
    sx2 += 1.2

# ===== 右下角版本信息 =====
add_text(slide, "数据截至：2026-05-08", 10.5, 5.35, 2.5, 0.22, 8, False, (130, 150, 180), PP_ALIGN.RIGHT, False)
add_text(slide, "Enterprise Asset & Security Management Platform v1.3", 0.5, 5.35, 9.5, 0.22, 8, False, (130, 150, 180), PP_ALIGN.LEFT, False)

# ===== 分层连接线（垂直虚线）=====
lines_x = [3.0, 4.7, 6.4, 8.1, 9.8]
for lx in lines_x:
    conn = slide.shapes.add_shape(1, Inches(lx), Inches(1.78), Inches(0.015), Inches(0.64))
    conn.fill.solid(); conn.fill.fore_color.rgb = rgb(180, 180, 200)
    conn.line.fill.background()

# 水平分隔线
sep1 = slide.shapes.add_shape(1, Inches(0.4), Inches(1.82), Inches(12.5), Inches(0.01))
sep1.fill.solid(); sep1.fill.fore_color.rgb = rgb(200, 200, 220); sep1.line.fill.background()
sep2 = slide.shapes.add_shape(1, Inches(0.4), Inches(2.42), Inches(12.5), Inches(0.01))
sep2.fill.solid(); sep2.fill.fore_color.rgb = rgb(200, 200, 220); sep2.line.fill.background()
sep3 = slide.shapes.add_shape(1, Inches(0.4), Inches(3.36), Inches(12.5), Inches(0.01))
sep3.fill.solid(); sep3.fill.fore_color.rgb = rgb(200, 200, 220); sep3.line.fill.background()

output = "d:\\资产管理\\项目架构图_专业版.pptx"
prs.save(output)
print(f"OK: {output}")