"""生成专业项目架构图PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_box(slide, text, left, top, width, height, bg_color, text_color, font_size=10, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*bg_color)
    shape.line.color.rgb = RGBColor(180, 180, 200)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*text_color)
    return shape

def add_label(slide, text, left, top, width, height, bg_color, text_color, font_size=10, bold=True):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*bg_color)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*text_color)
    return shape

def add_para(slide, text, left, top, width, height, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return tb

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# === 背景 ===
bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
bg.line.fill.background()

# === 标题栏 ===
tb = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(0.88))
tb.fill.solid()
tb.fill.fore_color.rgb = RGBColor(26, 54, 93)
tb.line.fill.background()

add_para(slide, "企业资产与信息安全综合管理平台 - 系统架构图", 0.5, 0.18, 12, 0.4, 20, (255, 255, 255), True)
add_para(slide, "B/S架构  ·  Flask 2.x + MySQL 8.0 + Ollama AI", 0.5, 0.58, 12, 0.3, 10, (160, 190, 230), False)

# === 用户层 ===
add_label(slide, "用户层", 0.4, 1.1, 1.2, 0.32, (26, 54, 93), (255, 255, 255), 11)

users = [("管理员", (255, 103, 0)), ("部门管理员", (0, 120, 212)),
         ("普通员工", (100, 100, 120)), ("外部API", (140, 60, 180))]
ux = 1.7
for u, c in users:
    add_box(slide, u, ux, 1.06, 1.3, 0.42, c, (255, 255, 255), 10, True)
    ux += 1.5

# 箭头1
arr = slide.shapes.add_shape(1, Inches(6.15), Inches(1.55), Inches(0.1), Inches(0.35))
arr.fill.solid(); arr.fill.fore_color.rgb = RGBColor(150, 150, 170); arr.line.fill.background()

# === 网关层 ===
add_label(slide, "安全网关层", 4.8, 1.95, 1.4, 0.3, (255, 103, 0), (255, 255, 255), 10)
add_box(slide, "CSRF防护 · API密钥认证 · JWT会话 · 限流熔断", 5.2, 2.3, 2.9, 0.42, (255, 248, 235), (180, 90, 0), 9, False)

# 箭头2
arr2 = slide.shapes.add_shape(1, Inches(6.15), Inches(2.78), Inches(0.1), Inches(0.35))
arr2.fill.solid(); arr2.fill.fore_color.rgb = RGBColor(150, 150, 170); arr2.line.fill.background()

# === 服务层 ===
add_label(slide, "应用服务层", 0.4, 3.15, 1.2, 0.3, (26, 54, 93), (255, 255, 255), 10)

services = [
    ("认证服务\n/auth", (255, 103, 0)),
    ("人员管理\n/personnel", (0, 120, 212)),
    ("IT资产管理\n/assets", (15, 148, 100)),
    ("安全管控\n/security", (200, 60, 60)),
    ("AI助手\n/ai_assistant", (130, 50, 190)),
    ("系统管理\n/admin", (100, 100, 100)),
    ("日志审计\n/audit", (60, 130, 170)),
]
sx = 1.7
for s, c in services:
    add_box(slide, s, sx, 3.12, 1.4, 0.62, c, (255, 255, 255), 9, True)
    sx += 1.5

# 箭头3
arr3 = slide.shapes.add_shape(1, Inches(6.15), Inches(3.8), Inches(0.1), Inches(0.35))
arr3.fill.solid(); arr3.fill.fore_color.rgb = RGBColor(150, 150, 170); arr3.line.fill.background()

# === 数据层 ===
add_label(slide, "数据层", 0.4, 4.18, 1.2, 0.3, (26, 54, 93), (255, 255, 255), 10)

# MySQL
add_box(slide, "🗄  MySQL 8.0\n11张核心业务表 · 索引优化 · 主从复制", 1.7, 4.15, 3.8, 0.52, (235, 242, 255), (0, 60, 130), 10, True)

# Redis
add_box(slide, "⚡  Redis 缓存\n会话存储 · 权限缓存 · 限流计数", 5.8, 4.15, 2.8, 0.52, (255, 240, 248), (160, 40, 80), 10, True)

# Ollama
add_box(slide, "🤖  Ollama AI\n意图识别 · SQL生成 · 自然语言回答", 8.9, 4.15, 3.8, 0.52, (245, 240, 255), (100, 40, 170), 10, True)

# === 底部统计栏 ===
ft = slide.shapes.add_shape(1, 0, Inches(5.95), Inches(13.333), Inches(1.55))
ft.fill.solid(); ft.fill.fore_color.rgb = RGBColor(246, 248, 252); ft.line.fill.background()

add_label(slide, "系统规模", 0.4, 6.05, 1.0, 0.28, (26, 54, 93), (255, 255, 255), 10)

stats = [
    ("👥 93,000+\n员工档案", (26, 54, 93), (255, 255, 255)),
    ("💻 1,000+\n办公电脑", (26, 54, 93), (255, 255, 255)),
    ("🔧 500+\n工控机设备", (26, 54, 93), (255, 255, 255)),
    ("🔐 7级\n角色权限", (26, 54, 93), (255, 255, 255)),
    ("🤖 AI\n智能问答", (26, 54, 93), (255, 255, 255)),
    ("📊 50+ RESTful\nAPI接口", (26, 54, 93), (255, 255, 255)),
    ("🛡️ CSRF+XSS+SQL\n全防护", (26, 54, 93), (255, 255, 255)),
]
tx = 0.4; ty = 6.38
for stat, bg_c, fg_c in stats:
    add_box(slide, stat, tx, ty, 1.55, 0.72, bg_c, fg_c, 9, True)
    tx += 1.65

add_para(slide, "数据截至：2026-05-08", 10.5, 6.9, 2.8, 0.25, 9, (160, 160, 180), False, PP_ALIGN.RIGHT)

output_path = "d:\\资产管理\\项目架构图.pptx"
prs.save(output_path)
print(f"架构图已保存: {output_path}")