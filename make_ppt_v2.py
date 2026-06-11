"""领导汇报PPT - 统一配色版"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_NAVY   = (26, 54, 93)
C_BLUE   = (0, 82, 170)
C_ORANGE = (230, 100, 0)
C_LIGHT  = (245, 247, 252)
C_LINE   = (200, 210, 230)
C_WHITE  = (255, 255, 255)
C_TEXT   = (60, 60, 80)
C_SUBTEXT= (120, 130, 160)

def R(*c): return RGBColor(*c)

def rect(sl, l, t, w, h, fill, line=None, lw=0.75):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = R(*fill)
    if line:
        s.line.color.rgb = R(*line); s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def tb(sl, text, l, t, w, h, sz=9, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.color.rgb = R(*color)
    return box

def tb_multi(sl, lines, l, t, w, h, sizes, colors, bolds, align=PP_ALIGN.LEFT):
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i,(line,sz,col,bold) in enumerate(zip(lines,sizes,colors,bolds)):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.size = Pt(sz); run.font.bold = bold
        run.font.color.rgb = R(*col)
    return box

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])

# 白色背景
rect(s, 0, 0, 13.333, 7.5, C_WHITE)

# ===== 顶部横幅 =====
rect(s, 0, 0, 13.333, 1.0, C_NAVY)
rect(s, 0, 0.95, 13.333, 0.05, C_ORANGE)

tb(s, "ENTERPRISE ASSET & SECURITY MANAGEMENT", 0.5, 0.1, 8, 0.3, 8, False, (160,190,230))
tb(s, "企业资产与信息安全综合管理平台", 0.5, 0.35, 10, 0.5, 20, True)
tb(s, "B/S三层架构  |  Flask + MySQL + Ollama AI  |  v1.3", 0.5, 0.78, 12, 0.25, 9, False, (140,170,210))
tb(s, "2026年5月8日", 11.2, 0.78, 1.9, 0.25, 9, False, (140,170,210), PP_ALIGN.RIGHT)

# ===== 第一行：三大成果 =====
y = 1.15
cards = [
    ("数据集中化", "93,000+", "员工档案统一管理，消除信息孤岛"),
    ("权限精细化", "7级RBAC", "最小粒度到人-系统-操作，部门数据隔离"),
    ("AI智能化", "自然语言查数据", "非技术人员自主查询，秒级响应"),
]
cx = 0.5
for title, num, desc in cards:
    rect(s, cx, y, 4.1, 0.95, C_NAVY)
    tb(s, title, cx+0.2, y+0.1, 3.7, 0.25, 9, True)
    tb(s, num, cx+0.2, y+0.38, 3.7, 0.3, 16, True)
    tb(s, desc, cx+0.2, y+0.68, 3.7, 0.22, 8, False, (180,210,240))
    cx += 4.25

# 箭头
rect(s, 6.5, y+0.95, 0.1, 0.2, C_LINE)
tb(s, "▼", 6.4, y+0.92, 0.3, 0.25, 12, False, (160,180,210))

# ===== 第二行：已完成 =====
y2 = y + 1.2
rect(s, 0.5, y2, 12.3, 0.3, C_LIGHT, C_BLUE, 0.5)
tb(s, "已上线运行  |  6大核心模块", 0.65, y2+0.04, 12, 0.24, 9, True, C_BLUE)

done = [
    ("人员管理系统", "9.3万员工档案"),
    ("IT资产管理", "1000+办公电脑+500+工控机"),
    ("安全管控平台", "涉密人/物/区域/文件"),
    ("AI智能问答", "意图识别+SQL生成+自然语言"),
    ("权限安全加固", "CSRF+XSS+密码策略"),
    ("界面统一优化", "小米设计风格"),
]
dx = 0.5
for name, stats in done:
    rect(s, dx, y2+0.4, 2.0, 0.65, C_BLUE)
    tb(s, name, dx+0.12, y2+0.46, 1.76, 0.28, 10, True)
    tb(s, stats, dx+0.12, y2+0.75, 1.76, 0.22, 7.5, False, (200,220,255))
    dx += 2.07

# 箭头
rect(s, 6.5, y2+1.05, 0.1, 0.2, C_LINE)
tb(s, "▼", 6.4, y2+1.02, 0.3, 0.25, 12, False, (160,180,210))

# ===== 第三行：收尾中 =====
y3 = y2 + 1.25
rect(s, 0.5, y3, 12.3, 0.3, C_LIGHT, C_ORANGE, 0.5)
tb(s, "收尾中  |  安全测试 + 全面回归测试 + 性能优化", 0.65, y3+0.04, 12, 0.24, 9, True, (180,80,0))

progs = [
    ("日志审计收尾", "操作留痕，合规审计追溯", 90),
    ("安全测试收尾", "代码审计+渗透测试", 90),
    ("全面回归测试", "功能+性能+安全验证", 75),
    ("性能优化", "索引+字段+导入效率", 60),
    ("文档完善", "用户手册+部署文档", 50),
]
px = 0.5
for name, desc, pct in progs:
    pw = 2.35 * pct / 100
    rect(s, px, y3+0.4, 2.35, 0.55, (248,250,255), C_LINE, 0.5)
    rect(s, px, y3+0.4, pw, 0.55, C_NAVY)
    tb(s, f"{pct}%", px+1.75, y3+0.44, 0.5, 0.28, 9, True, C_NAVY, PP_ALIGN.RIGHT)
    tb(s, name, px+0.12, y3+0.44, 1.6, 0.26, 9, True, C_NAVY)
    tb(s, desc, px+0.12, y3+0.7, 2.2, 0.22, 7.5, False, C_TEXT)
    px += 2.45

# 箭头
rect(s, 6.5, y3+1.0, 0.1, 0.2, C_LINE)
tb(s, "▼", 6.4, y3+0.97, 0.3, 0.25, 12, False, (160,180,210))

# ===== 第四行：里程碑 =====
y4 = y3 + 1.2
rect(s, 0.5, y4, 12.3, 0.3, C_LIGHT, C_BLUE, 0.5)
tb(s, "里程碑  |  上线倒计时", 0.65, y4+0.04, 5, 0.24, 9, True, C_BLUE)

ms = [
    ("5/14", "安全收尾完成"),
    ("5/21", "全面测试完成"),
    ("5/28", "优化文档完成"),
    ("6/7", "正式上线运行"),
]
mx = 0.7
for date, title in ms:
    rect(s, mx, y4+0.4, 2.85, 0.55, C_NAVY if "6/7" in date else (60,100,160))
    tb(s, date, mx+0.12, y4+0.44, 1.0, 0.28, 14, True)
    tb(s, title, mx+0.12, y4+0.72, 2.6, 0.22, 9, True, (180,210,240))
    mx += 3.05

# 里程碑箭头
for ax in [3.65, 6.7, 9.75]:
    rect(s, ax, y4+0.66, 0.2, 0.04, C_LINE)
    tb(s, "▶", ax-0.08, y4+0.58, 0.35, 0.2, 10, False, C_LINE)

# ===== 底部统计栏 =====
rect(s, 0, 5.75, 13.333, 1.2, C_NAVY)
tb(s, "系统规模", 0.5, 5.82, 1.5, 0.25, 9, True, (180,200,230))

bot = [
    ("93,000+", "员工档案"),
    ("1,000+", "办公电脑"),
    ("500+", "工控机"),
    ("200+", "部门"),
    ("7级", "角色权限"),
    ("CSRF+XSS+SQL", "安全防护"),
    ("AI", "智能问答"),
    ("50+", "API接口"),
    ("秒级", "查询响应"),
]
bx = 0.5
for num, unit in bot:
    rect(s, bx, 6.1, 1.2, 0.72, (40, 70, 110))
    tb(s, num, bx+0.06, 6.14, 1.08, 0.32, 11, True)
    tb(s, unit, bx+0.06, 6.46, 1.08, 0.3, 8, False, (150,180, 220), PP_ALIGN.CENTER)
    bx += 1.27

# 预计上线高亮
rect(s, 12.0, 5.75, 1.2, 1.2, C_ORANGE)
tb(s, "6/7", 12.06, 5.85, 1.0, 0.38, 16, True)
tb(s, "预计上线", 12.06, 6.2, 1.0, 0.28, 8, True)
tb(s, "v1.3  |  2026-05-08", 0.5, 6.88, 12, 0.22, 8, False, (130,150,180), PP_ALIGN.CENTER)

prs.save("d:\\资产管理\\领导汇报_统一配色.pptx")
print("OK")
