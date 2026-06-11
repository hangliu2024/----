"""领导汇报PPT - 自上而下结构"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def rgb(*c): return RGBColor(*c)

def add_rect(slide, l, t, w, h, fill, line=None, lw=0.75):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(*fill)
    if line:
        s.line.color.rgb = rgb(*line)
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def add_tb(slide, text, l, t, w, h, sz, bold=False, color=(255,255,255), align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(*color)
    return tb

def add_tb_multi(slide, lines, l, t, w, h, sizes, colors, bolds, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (line, sz, col, bold) in enumerate(zip(lines, sizes, colors, bolds)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = rgb(*col)
    return tb

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])

# 白色背景
add_rect(s, 0, 0, 13.333, 7.5, (255, 255, 255))

# 顶部深蓝横幅
add_rect(s, 0, 0, 13.333, 1.1, (26, 54, 93))
add_rect(s, 0, 1.06, 13.333, 0.05, (255, 103, 0))

add_tb(s, "ENTERPRISE ASSET & SECURITY MANAGEMENT", 0.5, 0.12, 10, 0.3, 8, False, (160, 190, 230))
add_tb(s, "企业资产与信息安全综合管理平台", 0.5, 0.4, 10, 0.55, 22, True)
add_tb(s, "B/S三层架构  |  Flask + MySQL + Ollama AI  |  v1.3", 0.5, 0.85, 10, 0.25, 9, False, (140, 170, 210))
add_tb(s, "汇报日期：2026年5月8日", 10.8, 0.4, 2.3, 0.3, 9, False, (140, 170, 210), PP_ALIGN.RIGHT)

# ===== 第一行：三大核心成果 =====
y1 = 1.25
cards = [
    ("数据集中化", "9.3万员工档案统一管理", "消除信息孤岛，实时掌握全集团人员信息", (255, 103, 0)),
    ("权限精细化", "最小粒度到人-系统-操作", "部门数据隔离，满足保密资质要求", (0, 112, 210)),
    ("AI智能化", "自然语言查数据，秒级响应", "非技术人员自主查询，解放IT生产力", (112, 48, 160)),
]
cx = 0.5
for title, num, desc, c in cards:
    add_rect(s, cx, y1, 3.9, 1.05, c)
    add_tb(s, title, cx + 0.15, y1 + 0.1, 3.6, 0.3, 11, True)
    add_tb(s, num, cx + 0.15, y1 + 0.4, 3.6, 0.3, 16, True, (255, 255, 255))
    add_tb(s, desc, cx + 0.15, y1 + 0.7, 3.6, 0.3, 8, False, (220, 220, 255))
    cx += 4.1

# 连接箭头（第一行 → 第二行）
ay = 2.38
add_rect(s, 6.5, ay, 0.12, 0.25, (200, 200, 220))
add_tb(s, "▼", 6.44, ay - 0.02, 0.25, 0.3, 14, False, (180, 180, 200))

# ===== 第二行：已完成模块 =====
y2 = 2.7
add_rect(s, 0.5, y2, 12.3, 0.32, (240, 247, 252), (0, 112, 210), 0.5)
add_tb(s, "已完成  |  6大核心模块已上线运行", 0.65, y2 + 0.04, 12, 0.26, 10, True, (0, 80, 160))

done_modules = [
    ("人员管理系统", "9.3万员工 · 80+字段", (52, 211, 153)),
    ("IT资产系统", "1000+办公电脑+500+工控机", (0, 112, 210)),
    ("安全管控平台", "涉密人/物/区域/文件", (200, 60, 60)),
    ("AI智能问答", "意图识别+SQL生成+自然语言", (112, 48, 160)),
    ("权限安全加固", "CSRF+XSS+密码策略", (255, 103, 0)),
    ("界面统一优化", "小米设计风格一致体验", (0, 128, 128)),
]
mx = 0.5
for name, stats, c in done_modules:
    add_rect(s, mx, y2 + 0.42, 2.0, 0.65, c)
    add_tb(s, name, mx + 0.1, y2 + 0.46, 1.8, 0.3, 10, True)
    add_tb(s, stats, mx + 0.1, y2 + 0.76, 1.8, 0.28, 7.5, False, (220, 220, 255))
    mx += 2.07

# 连接箭头
ay2 = y2 + 1.15
add_rect(s, 6.5, ay2, 0.12, 0.25, (200, 200, 220))
add_tb(s, "▼", 6.44, ay2 - 0.02, 0.25, 0.3, 14, False, (180, 180, 200))

# ===== 第三行：收尾进度 =====
y3 = 3.9
add_rect(s, 0.5, y3, 12.3, 0.32, (255, 248, 235), (255, 103, 0), 0.5)
add_tb(s, "收尾中  |  安全测试 + 全面回归测试 + 性能优化", 0.65, y3 + 0.04, 12, 0.26, 10, True, (180, 80, 0))

progress = [
    ("日志审计收尾", "操作留痕 · 合规审计追溯", 90, (52, 211, 153)),
    ("安全测试收尾", "代码审计 · 渗透测试", 90, (255, 193, 7)),
    ("全面回归测试", "功能验证 · 性能测试", 75, (0, 112, 210)),
    ("数据库优化", "索引 · 字段长度 · 导入效率", 60, (0, 128, 128)),
    ("文档完善", "用户手册 · 部署文档", 50, (80, 80, 80)),
]
px = 0.5
for name, desc, pct, c in progress:
    pw = 2.35 * pct / 100
    add_rect(s, px, y3 + 0.4, 2.35, 0.55, (248, 250, 252), (200, 200, 220), 0.5)
    if pct > 0:
        add_rect(s, px, y3 + 0.4, pw, 0.55, c)
    add_tb(s, f"{pct}%", px + 1.75, y3 + 0.44, 0.55, 0.3, 9, True, c if pct > 0 else (150, 150, 150), PP_ALIGN.RIGHT)
    add_tb(s, name, px + 0.1, y3 + 0.44, 1.6, 0.26, 9, True, (26, 54, 93))
    add_tb(s, desc, px + 0.1, y3 + 0.7, 2.2, 0.22, 7, False, (100, 100, 100))
    px += 2.45

# ===== 里程碑时间轴 =====
ay3 = 4.98
add_rect(s, 0.5, ay3, 12.3, 0.32, (235, 242, 255), (0, 80, 160), 0.5)
add_tb(s, "里程碑  |  上线倒计时", 0.65, ay3 + 0.04, 5, 0.26, 10, True, (0, 60, 130))

ms = [
    ("5/14", "安全收尾完成", (52, 211, 153)),
    ("5/21", "全面测试完成", (0, 112, 210)),
    ("5/28", "优化文档完成", (0, 112, 210)),
    ("6/7", "正式上线运行", (255, 103, 0)),
]
mx2 = 0.7
for date, title, c in ms:
    add_rect(s, mx2, ay3 + 0.42, 2.85, 0.58, c)
    add_tb(s, date, mx2 + 0.1, ay3 + 0.46, 0.8, 0.28, 14, True)
    add_tb(s, title, mx2 + 0.1, ay3 + 0.75, 2.65, 0.22, 9, True, (220, 220, 255))
    mx2 += 3.05

# 里程碑间的箭头
for ax in [3.65, 6.7, 9.75]:
    add_rect(s, ax, ay3 + 0.7, 0.2, 0.04, (200, 200, 220))

# ===== 底部统计数据栏 =====
add_rect(s, 0, 5.85, 13.333, 1.2, (26, 54, 93))
add_tb(s, "系统规模", 0.5, 5.92, 1.5, 0.26, 9, True, (180, 200, 230))

bottom_stats = [
    ("93,000+", "员工档案"),
    ("1,000+", "办公电脑"),
    ("500+", "工控机"),
    ("200+", "部门"),
    ("7级", "角色权限"),
    ("CSRF+XSS+SQL", "安全防护"),
    ("AI", "智能问答"),
    ("50+", "API接口"),
    ("B/S", "混合架构"),
    ("秒级", "查询响应"),
]
bsx = 0.5
for num, unit in bottom_stats:
    add_rect(s, bsx, 6.2, 1.2, 0.72, (40, 70, 110))
    add_tb(s, num, bsx + 0.05, 6.22, 1.1, 0.34, 11, True)
    add_tb(s, unit, bsx + 0.05, 6.56, 1.1, 0.32, 7.5, False, (150, 180, 220), PP_ALIGN.CENTER)
    bsx += 1.26

add_tb(s, "数据截至：2026-05-08", 10.5, 6.92, 2.5, 0.22, 8, False, (130, 150, 180), PP_ALIGN.RIGHT)
add_tb(s, "Enterprise Asset & Security Management Platform v1.3", 0.5, 6.92, 10, 0.22, 8, False, (130, 150, 180))

prs.save("d:\\资产管理\\领导汇报_一页流式版.pptx")
print("OK")
